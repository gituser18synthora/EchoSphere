"""
Document loading module for various file formats
Supports PDF, Word, Excel, PowerPoint, text, Markdown, and more
Integrated with cleaners.py and chunker.py for production RAG
Accepts local file paths and io.BytesIO streams (remote URL download was
removed in the EchoSphere port -- see resolve_shared_server_path).

NOTE: every function in this module is BLOCKING (synchronous file parsing,
PyMuPDF rendering and, via parsing.ocr, tesseract / sync OpenAI calls). They
must run in worker threads via asyncio.to_thread(...) or an executor -- never
directly on the event loop.

FIXES vs previous version:
  1. load_markdown() — no longer converts md → HTML → strip tags.
     Returns raw markdown so heading structure is preserved.
  2. load_document() — 'md'/'markdown' files get file_type='markdown'
     (not 'text') so chunker uses header-aware splitting.
  3. preprocess_for_embedding() called with is_markdown=True for .md
     files so cleaners.py skips remove_markdown_formatting.
"""

import io
import re
import statistics
from collections import Counter
from typing import Union, Optional, List, Dict
from pathlib import Path
import mimetypes
import logging
import fitz
import json

from backend.knowledge.kconfig import settings

logger = logging.getLogger(__name__)

# PDF handling
try:
    import PyPDF2
    HAS_PYPDF2 = True
except ImportError:
    HAS_PYPDF2 = False

try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False

# Office documents
try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# Import cleaners and chunker
try:
    from backend.knowledge.parsing.cleaners import preprocess_for_embedding
    HAS_CLEANERS = True
except ImportError:
    HAS_CLEANERS = False
    logger.warning("cleaners.py not found, text cleaning disabled")

try:
    from backend.knowledge.chunking.chunker import create_chunks, CHUNK_CONFIGS
    HAS_CHUNKER = True
except ImportError:
    HAS_CHUNKER = False
    logger.warning("chunker.py not found, chunking disabled")


def resolve_shared_server_path(file_path: str) -> str:
    """
    No-op passthrough kept for API compatibility with the KMRAG source.

    The original implementation downloaded HTTP/HTTPS URLs to a temp file
    (with authenticated remote download). EchoSphere only ingests local files
    (or io.BytesIO streams, which never reach this function), so the path
    string is returned unchanged. Existence checking happens in load_document.
    """
    return str(file_path)


def load_text_file(file_path: Union[str, Path, io.BytesIO]) -> str:
    """Load plain text file"""
    try:
        if isinstance(file_path, io.BytesIO):
            content = file_path.read()
            return content.decode('utf-8', errors='ignore')
        else:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
    except Exception as e:
        logger.error(f"Failed to load text file: {e}")
        return ""


def load_pdf_pypdf2(file_path: Union[str, Path, io.BytesIO]) -> str:
    """Load PDF using PyPDF2"""
    if not HAS_PYPDF2:
        raise ImportError("PyPDF2 is not installed. Install with: pip install PyPDF2")
    try:
        text_content = []
        if isinstance(file_path, io.BytesIO):
            reader = PyPDF2.PdfReader(file_path)
        else:
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
        for page_num, page in enumerate(reader.pages, 1):
            try:
                text = page.extract_text()
                if text:
                    text_content.append(text)
            except Exception as e:
                logger.warning(f"Failed to extract page {page_num}: {e}")
        return "\n\n".join(text_content)
    except Exception as e:
        logger.error(f"Failed to load PDF with PyPDF2: {e}")
        return ""


def load_pdf_pdfplumber(file_path: Union[str, Path, io.BytesIO]) -> str:
    """Load PDF using pdfplumber (better table extraction)"""
    if not HAS_PDFPLUMBER:
        raise ImportError("pdfplumber is not installed. Install with: pip install pdfplumber")
    try:
        text_content = []
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                try:
                    text = page.extract_text()
                    if text:
                        text_content.append(text)

                    tables = page.extract_tables()
                    for table in tables:
                        table_text = "\n".join([" | ".join([str(cell) for cell in row]) for row in table])
                        text_content.append(f"\n[TABLE]\n{table_text}\n[/TABLE]\n")
                except Exception as e:
                    logger.warning(f"Failed to extract page {page_num}: {e}")
        return "\n\n".join(text_content)
    except Exception as e:
        logger.error(f"Failed to load PDF with pdfplumber: {e}")
        return ""



def table_to_html(table) -> str:
    html = ["<table border='1' style='border-collapse: collapse;'>"]
    rows = table.extract()

    for i, row in enumerate(rows):
        html.append("<tr>")
        for cell in row:
            tag = "th" if i == 0 else "td"
            text = (cell or "").replace("\n", " ").strip()
            html.append(f"<{tag}>{text}</{tag}>")
        html.append("</tr>")

    html.append("</table>")
    return "\n".join(html)


def table_to_html_word(table) -> str:
    """Convert Word table to HTML"""
    html = ["<table border='1' style='border-collapse: collapse;'>"]

    for i, row in enumerate(table.rows):
        html.append("<tr>")
        for cell in row.cells:
            tag = "th" if i == 0 else "td"
            text = cell.text.strip().replace("\n", " ")
            html.append(f"<{tag}>{text}</{tag}>")
        html.append("</tr>")

    html.append("</table>")
    return "\n".join(html)


def overlaps(b1, b2):
    return not (
        b1[2] < b2[0] or
        b1[0] > b2[2] or
        b1[3] < b2[1] or
        b1[1] > b2[3]
    )


def _pymupdf_page_texts(file_path: Union[str, Path, io.BytesIO]) -> List[str]:
    """
    Extract inline text + tables in reading order, returning ONE string per page
    (index 0 = page 1). Keeping the per-page split lets the caller attach
    page_number to each chunk; load_pdf_pymupdf() just joins these for callers
    that want a flat string.
    """
    if isinstance(file_path, io.BytesIO):
        doc = fitz.open(stream=file_path.read(), filetype="pdf")
    else:
        doc = fitz.open(str(file_path))

    page_texts: List[str] = []
    try:
        for page_index, page in enumerate(doc):
            page_number = page_index + 1

            elements = []

            tables = []
            try:
                tables = list(page.find_tables())
            except Exception:
                pass

            table_boxes = []
            rejected_tables = 0
            for table in tables:
                try:
                    rows = table.extract()
                except Exception:
                    rows = []
                # Reject false-positive tables (1-row / 1-column detections that
                # find_tables() emits on whitespace-heavy prose). If trusted, every
                # text block overlapping their bbox is dropped below — and because
                # get_text("blocks") returns coarse multi-line blocks, a single
                # phantom table can delete a whole paragraph of real content. See
                # _is_real_table for the generic test.
                if not _is_real_table(rows):
                    rejected_tables += 1
                    continue
                bbox = table.bbox
                table_boxes.append(bbox)
                elements.append({
                    "type": "table",
                    "bbox": bbox,
                    "content": table_to_html(table)
                })

            blocks = page.get_text("blocks")

            kept_blocks = 0
            dropped_blocks = 0
            for block in blocks:
                x0, y0, x1, y1, text, *_ = block
                block_box = (x0, y0, x1, y1)

                if any(overlaps(block_box, tbl_box) for tbl_box in table_boxes):
                    dropped_blocks += 1
                    continue

                clean_text = text.strip()
                if clean_text:
                    kept_blocks += 1
                    elements.append({
                        "type": "text",
                        "bbox": block_box,
                        "content": clean_text
                    })

            elements.sort(key=lambda e: (round(e["bbox"][1], 1), round(e["bbox"][0], 1)))

            parts = []
            page_text_chars = 0
            for el in elements:
                if el["type"] == "table":
                    parts.append(f"<doc-table>\n{el['content']}\n</doc-table>")
                else:
                    parts.append(el["content"])
                    page_text_chars += len(el["content"])

            page_texts.append("\n\n".join(parts))

            logger.info(
                "PDF(flat) page %d: tables_found=%d real_tables=%d rejected_tables=%d "
                "blocks_kept=%d blocks_dropped=%d text_chars=%d",
                page_number, len(tables), len(table_boxes), rejected_tables,
                kept_blocks, dropped_blocks, page_text_chars,
            )

    finally:
        doc.close()

    return page_texts


def load_pdf_pymupdf(file_path: Union[str, Path, io.BytesIO]) -> str:
    """Inline text + tables in correct reading order (flat string, all pages)."""
    return "\n\n".join(t for t in _pymupdf_page_texts(file_path) if t.strip())


def load_pdf_pages(file_path: Union[str, Path, io.BytesIO]) -> List[Dict]:
    """
    Page-aware PDF text for the standard (non-layout-aware) path.

    Returns [{"page_number": int, "text": str}, ...] so load_document can chunk
    each page independently and stamp page_number onto every chunk. Returns []
    on failure so the caller can fall back to flat whole-document extraction.
    """
    try:
        texts = _pymupdf_page_texts(file_path)
    except Exception as e:
        logger.warning(f"Page-aware PDF extraction failed: {e}")
        return []
    return [
        {"page_number": i + 1, "text": t}
        for i, t in enumerate(texts)
        if t and t.strip()
    ]

def load_pdf(file_path: Union[str, Path, io.BytesIO], prefer_pdfplumber: bool = True) -> str:
    """Load PDF with automatic fallback and debugging"""
    logger.info(f"Attempting to load PDF: {file_path}")

    try:
        text = load_pdf_pymupdf(file_path)
        if text and text.strip():
            return text
    except Exception as e:
        logger.warning(f"pdf extraction failed: {str(e)}")


# ==================== LAYOUT-AWARE PDF EXTRACTION ====================
# Page-structured extraction for PPT-style/table-heavy PDFs. Separate from
# load_pdf_pymupdf() above (which stays untouched) so existing callers keep
# getting a flat string. This path is only used when explicitly requested
# (see load_document()'s ENABLE_LAYOUT_AWARE_CHUNKING branch).

_BULLET_PATTERN = re.compile(r'^\s*([•▪‣◦∙·–\-\*]|\d+[\.\)]|[a-zA-Z][\.\)])\s+')


def _is_bullet_line(text: str) -> bool:
    return bool(_BULLET_PATTERN.match(text.strip()))


# A real heading/section title must contain actual words, not just a number,
# a percentage, or symbols. Without this, a large-font "0%" / "100%" / "—" in a
# promo graphic gets misdetected as a heading and becomes the section name
# (the "Section: 0%" bug). Matches Latin + Latin-1 + Devanagari letters.
_WORDISH = re.compile(r"[A-Za-zÀ-ɏऀ-ॿ]")


def _is_heading_textlike(text: str) -> bool:
    """True only if the candidate heading has >=2 alphabetic characters."""
    return len(_WORDISH.findall(text or "")) >= 2


def _is_real_table(rows: List[List]) -> bool:
    """
    Reject PyMuPDF false-positive tables before they're trusted.

    `find_tables()` regularly flags a stray aligned phrase on a whitespace-heavy
    prose page as a 1-row / 1-column "table" (e.g. a "Check | Status" label). If
    that phantom is trusted, two bad things happen: every text line overlapping
    its bbox is dropped from the page (silent data loss), and the table itself
    yields no rows so nothing replaces it. A genuine grid needs at least two
    non-empty rows AND at least two columns — anything less is treated as plain
    text instead.
    """
    non_empty_rows = [
        r for r in (rows or [])
        if any(c is not None and str(c).strip() for c in r)
    ]
    if len(non_empty_rows) < 2:
        return False
    max_cols = max(len(r) for r in non_empty_rows)
    return max_cols >= 2


def _normalize_for_repeat_check(text: str) -> str:
    """Collapse digit runs (page numbers, dates) so 'Page 2'/'Page 3' compare equal."""
    return re.sub(r"\d+", "#", text.strip().lower())


def _strip_repeated_header_footer(pages: List[Dict]) -> None:
    """Drop lines repeated (modulo digits) as the first/last block on most pages (running headers/footers)."""
    if len(pages) < 4:
        return

    def _edge_texts(page):
        blocks = [b for b in page["blocks"] if b["type"] != "table"]
        texts = []
        if blocks:
            texts.append(blocks[0]["content"].strip())
            if len(blocks) > 1:
                texts.append(blocks[-1]["content"].strip())
        return texts

    counter = Counter()
    for page in pages:
        for t in {_normalize_for_repeat_check(t) for t in _edge_texts(page) if t}:
            counter[t] += 1

    threshold = max(3, int(len(pages) * 0.6))
    repeated = {t for t, c in counter.items() if c >= threshold}
    if not repeated:
        return

    for page in pages:
        page["blocks"] = [
            b for b in page["blocks"]
            if not (b["type"] != "table" and _normalize_for_repeat_check(b["content"]) in repeated)
        ]


# ==================== OCR FALLBACK (image / scanned PDF text) ====================
# Many PDFs carry text and tables inside images (screenshots, scanned pages,
# embedded image-tables, forms, charts) that the normal text layer does not
# expose. When ENABLE_OCR_FALLBACK is on and a page is weak on native text OR
# image-heavy, we OCR it (local tesseract first, optional GPT-vision escalation)
# and merge only the *new* content into the page so we never duplicate text into
# extra chunks. All OCR mechanics live in ingestion/ocr.py.

from backend.knowledge.parsing.ocr import (
    OcrBudget,
    analyze_page_images,
    ocr_enabled,
    run_page_ocr,
)


def extract_pdf_pages_structured(file_path: Union[str, Path, io.BytesIO]) -> List[Dict]:
    """
    Layout-aware PDF extraction for PPT-style/table-heavy PDFs.

    Returns one dict per page:
        {"page_number": int, "blocks": [{"type": "heading"|"text"|"table", "content": str, "rows": [...]}]}

    - Headings are detected by font size relative to the page's median font size.
    - Consecutive bullet lines are grouped into one block.
    - Tables are kept atomic (never split) and carry raw `rows` for NL conversion.
    - Lines repeated verbatim across most pages (running headers/footers) are dropped.
    """
    if isinstance(file_path, io.BytesIO):
        doc = fitz.open(stream=file_path.read(), filetype="pdf")
    else:
        doc = fitz.open(str(file_path))

    pages_out: List[Dict] = []
    ocr_on = ocr_enabled()
    ocr_budget = OcrBudget()  # per-document GPT-vision budget
    try:
        for page_index, page in enumerate(doc):
            page_number = page_index + 1

            tables = []
            try:
                tables = page.find_tables()
            except Exception:
                pass

            table_boxes = []
            table_elements = []
            for table in tables:
                bbox = table.bbox
                try:
                    rows = table.extract()
                except Exception:
                    rows = []
                # Skip false-positive tables: don't register their bbox (so the
                # prose lines they cover stay as text) and don't emit them.
                if not _is_real_table(rows):
                    continue
                table_boxes.append(bbox)
                table_elements.append({
                    "type": "table",
                    "bbox": bbox,
                    "content": table_to_html(table),
                    "rows": rows,
                })
            native_table_count = len(table_elements)
            # Cell text of natively-extracted tables, fed to OCR de-dup so a page
            # whose table is ALSO inside an image doesn't get its cells re-emitted
            # as a free-text OCR chunk that duplicates the table_summary chunks.
            native_table_text = "\n".join(
                str(cell).strip()
                for el in table_elements
                for row in (el.get("rows") or [])
                for cell in row
                if cell and str(cell).strip()
            )

            text_dict = page.get_text("dict")
            line_records = []
            sizes = []

            for block in text_dict.get("blocks", []):
                if block.get("type") != 0:
                    continue
                for line in block.get("lines", []):
                    spans = line.get("spans", [])
                    if not spans:
                        continue
                    line_text = "".join(s.get("text", "") for s in spans).strip()
                    if not line_text:
                        continue
                    bbox = line.get("bbox", (0, 0, 0, 0))
                    if any(overlaps(bbox, tbox) for tbox in table_boxes):
                        continue
                    max_size = max(s.get("size", 0) for s in spans)
                    sizes.append(max_size)
                    line_records.append({"text": line_text, "bbox": bbox, "size": max_size})

            median_size = statistics.median(sizes) if sizes else 0

            elements = list(table_elements)
            bullet_buffer: List[str] = []
            bullet_bbox = None
            last_heading_idx = None

            def flush_bullets():
                nonlocal bullet_buffer, bullet_bbox
                if bullet_buffer:
                    elements.append({
                        "type": "text",
                        "bbox": bullet_bbox,
                        "content": "\n".join(bullet_buffer),
                    })
                    bullet_buffer = []
                    bullet_bbox = None

            for rec in line_records:
                text, bbox, size = rec["text"], rec["bbox"], rec["size"]
                if _is_bullet_line(text):
                    bullet_buffer.append(text if text.lstrip().startswith(("-", "*")) else f"- {text}")
                    bullet_bbox = bbox
                    last_heading_idx = None
                    continue

                flush_bullets()
                is_heading = (
                    median_size > 0
                    and size >= median_size * 1.15
                    and len(text) <= 120
                    and _is_heading_textlike(text)
                )

                # Merge consecutive heading lines, but only if they're actually
                # adjacent on the page (a wrapped title) — not just adjacent in
                # the filtered line stream, which can also happen when a table
                # sitting between two real headings has its text excluded above.
                if is_heading and last_heading_idx is not None:
                    prev = elements[last_heading_idx]
                    vertical_gap = bbox[1] - prev["bbox"][3]
                    if vertical_gap < max(size, prev.get("_size", size)) * 1.8:
                        prev["content"] += " " + text
                        prev["bbox"] = (
                            min(prev["bbox"][0], bbox[0]),
                            min(prev["bbox"][1], bbox[1]),
                            max(prev["bbox"][2], bbox[2]),
                            max(prev["bbox"][3], bbox[3]),
                        )
                        prev["_size"] = size
                        continue

                elements.append({
                    "type": "heading" if is_heading else "text",
                    "bbox": bbox,
                    "content": text,
                    "_size": size,
                })
                last_heading_idx = (len(elements) - 1) if is_heading else None

            flush_bullets()

            elements.sort(key=lambda e: (round(e["bbox"][1], 1), round(e["bbox"][0], 1)))
            for el in elements:
                el.pop("bbox", None)
                el.pop("_size", None)

            # OCR fallback: recover text and tables that live inside images. This
            # fires when EITHER the native text layer is weak (scanned/empty page)
            # OR the page is image-heavy (a screenshot/scanned table that still
            # carries a small caption/title in its text layer). The second
            # condition is what lets a page with some normal text — but its real
            # data inside an image — still get OCR'd, instead of being skipped just
            # because text > threshold. Provider selection + table extraction +
            # de-duplication all happen in ingestion/ocr.py.
            native_text = "\n".join(e["content"] for e in elements if e["type"] != "table")
            native_len = len(native_text.strip())
            image_info = analyze_page_images(page)

            ocr_summary = {
                "used": False, "provider": None, "from_image": False,
                "table_detected": False, "confidence": None, "language": None,
            }
            ocr_triggered = False
            ocr_text_len = 0
            ocr_tables = 0
            if ocr_on:
                weak_text = (
                    settings.OCR_PROCESS_SCANNED_PAGES
                    and native_len < settings.OCR_MIN_PAGE_CHARS
                )
                image_heavy = (
                    settings.OCR_PROCESS_EMBEDDED_IMAGES
                    and image_info.has_meaningful_image
                    and image_info.image_area_ratio >= settings.OCR_IMAGE_AREA_RATIO
                )
                if weak_text or image_heavy:
                    ocr_triggered = True
                    ocr_result = run_page_ocr(
                        page,
                        native_text=native_text,
                        page_number=page_number,
                        image_info=image_info,
                        native_table_count=native_table_count,
                        budget=ocr_budget,
                        native_table_text=native_table_text,
                    )
                    if ocr_result and ocr_result.used:
                        for ocr_block in ocr_result.blocks:
                            block = ocr_block.as_block()
                            elements.append(block)
                            if block["type"] == "table":
                                ocr_tables += 1
                            else:
                                ocr_text_len += len(block["content"])
                        ocr_summary = ocr_result.page_summary()
                    logger.info(
                        f"OCR page {page_number}: triggered "
                        f"(weak_text={weak_text}, image_heavy={image_heavy}, "
                        f"area_ratio={image_info.image_area_ratio:.2f}) -> provider="
                        f"{ocr_summary['provider']}, text={ocr_text_len} chars, "
                        f"tables={ocr_tables}"
                    )

            merged_len = sum(len(e["content"]) for e in elements if e["type"] != "table")
            logger.info(
                f"PDF page {page_number}: native_text={native_len} chars, "
                f"images={image_info.num_images} (meaningful={image_info.num_meaningful_images}, "
                f"area={image_info.image_area_ratio:.2f}), ocr_triggered={ocr_triggered}, "
                f"ocr_text={ocr_text_len} chars, ocr_tables={ocr_tables}, "
                f"merged_text={merged_len} chars, blocks={len(elements)}"
            )

            pages_out.append({
                "page_number": page_number,
                "blocks": elements,
                "ocr": ocr_summary,
            })
    finally:
        doc.close()

    _strip_repeated_header_footer(pages_out)
    return pages_out



def iter_block_items(doc):
    """
    Yield paragraphs and tables in document order
    """
    from docx.oxml.ns import qn

    for child in doc.element.body:
        if child.tag.endswith('p'):
            yield docx.text.paragraph.Paragraph(child, doc)
        elif child.tag.endswith('tbl'):
            yield docx.table.Table(child, doc)


def load_word_document(file_path: Union[str, Path, io.BytesIO]) -> str:
    """Extract text + inline tables (as HTML) in correct order"""

    try:
        if isinstance(file_path, io.BytesIO):
            doc = docx.Document(file_path)
        else:
            doc = docx.Document(file_path)

        output = []

        for block in iter_block_items(doc):

            if isinstance(block, docx.text.paragraph.Paragraph):
                text = block.text.strip()
                if not text:
                    continue
                if block.style.name.startswith("Heading"):
                    output.append(f"<h>{text}</h>")
                else:
                    output.append(text)

            elif isinstance(block, docx.table.Table):
                html_table = table_to_html_word(block)

                output.append(
                    f"<doc-table>\n{html_table}\n</doc-table>"
                )

        return "\n\n".join(output)

    except Exception as e:
        logger.error(f"Failed to load Word document: {e}")
        return ""



def load_excel_file(file_path: Union[str, Path, io.BytesIO]) -> dict:
    """Load Excel into structured dict with headers"""

    wb = openpyxl.load_workbook(file_path)
    excel_data = {}

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        rows = list(sheet.iter_rows(values_only=True))

        if not rows:
            continue

        # `if c` / `if cell` truthiness silently dropped every 0, 0.0 and False
        # cell (and turned a 0 header into col_N) — compare against None so
        # numeric zeros and booleans survive into the KB.
        headers = [
            str(c).strip() if c is not None and str(c).strip() else f"col_{i}"
            for i, c in enumerate(rows[0])
        ]

        sheet_data = []

        for row_idx, row in enumerate(rows[1:], start=1):
            values = [str(cell).strip() if cell is not None else "" for cell in row]

            if not any(values):
                continue

            row_dict = {
                headers[i]: values[i] if i < len(values) else ""
                for i in range(len(headers))
            }

            sheet_data.append({
                "row_index": row_idx,
                "data": row_dict
            })

        if sheet_data:
            excel_data[sheet_name] = {
                "headers": headers,
                "rows": sheet_data
            }

    return excel_data


def load_powerpoint(file_path: Union[str, Path, io.BytesIO]) -> str:
    """Load PowerPoint presentation (.pptx)"""
    if not HAS_PPTX:
        raise ImportError("python-pptx is not installed. Install with: pip install python-pptx")
    try:
        if isinstance(file_path, io.BytesIO):
            prs = Presentation(file_path)
        else:
            prs = Presentation(file_path)
        text_content = []

        for slide_num, slide in enumerate(prs.slides, 1):
            text_content.append(f"\n[SLIDE {slide_num}]\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_content.append(shape.text)
            text_content.append(f"[/SLIDE]\n")

        return "\n\n".join(text_content)
    except Exception as e:
        logger.error(f"Failed to load PowerPoint: {e}")
        return ""


def _table_to_html_pptx(table) -> str:
    """Convert a python-pptx table to HTML (mirrors table_to_html_word above)."""
    html = ["<table border='1' style='border-collapse: collapse;'>"]
    for i, row in enumerate(table.rows):
        html.append("<tr>")
        for cell in row.cells:
            tag = "th" if i == 0 else "td"
            text = cell.text.strip().replace("\n", " ")
            html.append(f"<{tag}>{text}</{tag}>")
        html.append("</tr>")
    html.append("</table>")
    return "\n".join(html)


def extract_pptx_slides_structured(file_path: Union[str, Path, io.BytesIO]) -> List[Dict]:
    """
    Layout-aware PPTX extraction: returns one dict per slide:
        {"page_number": int, "blocks": [{"type": "heading"|"text"|"table", "content": str, "rows": [...]}]}

    Title placeholders become headings; indented paragraph levels become
    grouped bullet blocks; tables are kept atomic with raw rows for NL conversion.
    """
    if not HAS_PPTX:
        raise ImportError("python-pptx is not installed. Install with: pip install python-pptx")

    prs = Presentation(file_path)
    slides_out: List[Dict] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        blocks: List[Dict] = []
        bullet_buffer: List[str] = []

        def flush_bullets():
            nonlocal bullet_buffer
            if bullet_buffer:
                blocks.append({"type": "text", "content": "\n".join(bullet_buffer)})
                bullet_buffer = []

        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                flush_bullets()
                try:
                    rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                except Exception:
                    rows = []
                blocks.append({
                    "type": "table",
                    "content": _table_to_html_pptx(shape.table),
                    "rows": rows,
                })
                continue

            if not getattr(shape, "has_text_frame", False):
                continue

            is_title = False
            try:
                if shape.is_placeholder and shape.placeholder_format.type is not None:
                    is_title = shape.placeholder_format.type in (
                        PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE
                    )
            except Exception:
                is_title = False

            full_text = shape.text_frame.text.strip()
            if not full_text:
                continue

            if is_title:
                flush_bullets()
                blocks.append({"type": "heading", "content": full_text})
                continue

            for para in shape.text_frame.paragraphs:
                para_text = para.text.strip()
                if not para_text:
                    continue
                if (para.level or 0) > 0 or _is_bullet_line(para_text):
                    bullet_buffer.append(
                        para_text if para_text.startswith(("-", "*")) else f"- {para_text}"
                    )
                else:
                    flush_bullets()
                    blocks.append({"type": "text", "content": para_text})

        flush_bullets()
        slides_out.append({"page_number": slide_num, "blocks": blocks})

    return slides_out


def load_markdown(file_path: Union[str, Path, io.BytesIO]) -> str:
    """
    Load Markdown file — returns RAW markdown content.

    FIX: Previous version converted md → HTML → stripped tags,
    which destroyed all heading/table structure before chunking.
    Now we return raw .md text so chunker.py can split on ## headers.
    The markdown library import is intentionally not used here anymore.
    """
    try:
        if isinstance(file_path, io.BytesIO):
            return file_path.read().decode('utf-8', errors='ignore')
        else:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
    except Exception as e:
        logger.error(f"Failed to load Markdown: {e}")
        return ""


# Structure counters for raw markdown (run BEFORE cleaning so fences/tables
# are still intact). Setext headings and reference-style links are not counted.
_MD_HEADING_RE = re.compile(r"^#{1,6}\s+\S", re.MULTILINE)
_MD_LINK_RE = re.compile(r"(?<!\!)\[[^\]]*\]\([^)]+\)")
_MD_FENCE_RE = re.compile(r"^\s*(?:```|~~~)", re.MULTILINE)
_MD_TABLE_SEP_RE = re.compile(r"^(?=[^\n]*\|)\s*\|?[\s:|-]*-{3,}[\s:|-]*$", re.MULTILINE)


def markdown_content_stats(content: str) -> Dict:
    """Document-level structure stats for a Markdown file, stamped onto every
    chunk's metadata so extraction_metadata carries more than OCR provenance."""
    return {
        "heading_count": len(_MD_HEADING_RE.findall(content)),
        "link_count": len(_MD_LINK_RE.findall(content)),
        "code_block_count": len(_MD_FENCE_RE.findall(content)) // 2,
        "table_count": len(_MD_TABLE_SEP_RE.findall(content)),
    }


def load_csv(file_path: Union[str, Path, io.BytesIO]) -> dict:
    """
    Load CSV into the same structured dict shape as load_excel_file, so CSV
    goes through the header-aware row chunker (chunk_excel_smart) instead of
    flat text splitting. Flat splitting lost the header row after the first
    chunk — later chunks were bare `val | val | val` rows with no column names,
    useless to the LLM. The first row is treated as the header row.
    """
    try:
        import csv
        if isinstance(file_path, io.BytesIO):
            content = file_path.read().decode('utf-8', errors='ignore')
            rows = list(csv.reader(content.splitlines()))
            table_name = "CSV"
        else:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                rows = list(csv.reader(f))
            table_name = Path(file_path).stem or "CSV"

        rows = [r for r in rows if any(str(c).strip() for c in r)]
        if not rows:
            return {}

        headers = [
            str(c).strip() if c is not None and str(c).strip() else f"col_{i}"
            for i, c in enumerate(rows[0])
        ]

        sheet_data = []
        for row_idx, row in enumerate(rows[1:], start=1):
            values = [str(cell).strip() if cell is not None else "" for cell in row]
            if not any(values):
                continue
            row_dict = {
                headers[i]: values[i] if i < len(values) else ""
                for i in range(len(headers))
            }
            sheet_data.append({"row_index": row_idx, "data": row_dict})

        if not sheet_data:
            return {}

        return {table_name: {"headers": headers, "rows": sheet_data}}
    except Exception as e:
        logger.error(f"Failed to load CSV: {e}")
        return {}


def load_json(file_path: Union[str, Path, io.BytesIO]) -> str:
    """Load JSON file"""
    try:
        if isinstance(file_path, io.BytesIO):
            content = file_path.read().decode('utf-8', errors='ignore')
            data = json.loads(content)
        else:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)

        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception as e:
        logger.error(f"Failed to load JSON: {e}")
        return ""


def get_file_extension(file_path: Union[str, Path, io.BytesIO], filename: str = None) -> str:
    """Get file extension"""
    if isinstance(file_path, io.BytesIO):
        if filename:
            return Path(filename).suffix.lower().lstrip('.')
        return ""
    else:
        return Path(file_path).suffix.lower().lstrip('.')


def _try_layout_aware_chunking(file_path_obj, ext: str, source_label: str) -> Optional[List[Dict]]:
    """
    Attempt page/slide-structured extraction + structured chunking for
    PDF/PPT-style documents. Returns None (never raises) when the document
    doesn't qualify (e.g. a plain prose PDF) or extraction fails, so the
    caller falls back to the standard flat-text chunking path unchanged.
    """
    try:
        from backend.knowledge.parsing.doc_classifier import detect_document_type
        from backend.knowledge.chunking.structured_chunker import build_structured_chunks, build_embedding_text

        if ext == "pdf":
            pages = extract_pdf_pages_structured(file_path_obj)
        else:
            pages = extract_pptx_slides_structured(file_path_obj)

        if not pages:
            return None

        doc_type = detect_document_type(ext, pages)

        # A scanned / image-based PDF often classifies as plain prose
        # (pdf_text_document) yet its real content was just recovered by OCR.
        # Keep those on the structured path so the OCR'd text/tables (and their
        # extraction metadata) survive into chunks instead of being dropped by
        # the flat fallback, which never OCRs.
        ocr_present = any((p.get("ocr") or {}).get("used") for p in pages)

        if doc_type not in ("pdf_ppt_style", "table_heavy_pdf") and not ocr_present:
            return None  # plain prose PDF with no image content -> standard path

        document_name = Path(source_label).name

        strategy = "layout_aware"
        structured_chunks = build_structured_chunks(pages, file_name=source_label)
        if not structured_chunks:
            return None

        logger.info(f"Layout-aware chunking: {len(structured_chunks)} chunks "
                    f"(doc_type={doc_type}, strategy={strategy})")

        source_path = str(file_path_obj) if not isinstance(file_path_obj, io.BytesIO) else None

        results = []
        for idx, chunk in enumerate(structured_chunks):
            extraction = chunk.get("extraction") or {}
            results.append({
                "page_content": chunk["content"],
                "metadata": {
                    "file_type": ext,
                    "chunk_index": idx,
                    "chunk_strategy": strategy,
                    "doc_type": doc_type,
                    "source": source_label,
                    "file_path": source_path,
                    "file_name": document_name,
                    "page_number": chunk["page_number"],
                    "section": chunk["section"],
                    "topic": chunk["topic"],
                    "chunk_type": chunk["chunk_type"],
                    "keywords": chunk["keywords"],
                    "language": extraction.get("language"),
                    # OCR / image extraction provenance (see ingestion/ocr.py)
                    "extraction_method": extraction.get("extraction_method", "native_text"),
                    "ocr_used": extraction.get("ocr_used", False),
                    "ocr_provider": extraction.get("ocr_provider"),
                    "from_image": extraction.get("from_image", False),
                    "table_detected": extraction.get("table_detected", False),
                    "table_index": extraction.get("table_index"),
                    "confidence": extraction.get("confidence"),
                    "embedding_text": chunk.get("embedding_text") or build_embedding_text(
                        document_name=document_name,
                        section=chunk["section"],
                        topic=chunk["topic"],
                        chunk_type=chunk["chunk_type"],
                        keywords=chunk["keywords"],
                        content=chunk["content"],
                    ),
                },
            })
        return results
    except Exception as e:
        logger.warning(f"Layout-aware chunking failed, falling back: {e}")
        return None


def _chunk_pdf_pages_with_numbers(
    page_docs: List[Dict],
    filename: Optional[str],
    file_path,
    file_path_obj,
    enable_cleaning: bool,
    chunk_size: Optional[int],
    chunk_overlap: Optional[int],
) -> List[Dict]:
    """
    Chunk each PDF page independently (standard path) and stamp page_number onto
    every resulting chunk. Cleaning is applied per page so behaviour matches the
    flat path; only the page boundary is preserved so chunks remain page-scoped.
    """
    std_source_path = str(file_path_obj) if not isinstance(file_path_obj, io.BytesIO) else None
    out: List[Dict] = []
    global_idx = 0

    for page_doc in page_docs:
        page_number = page_doc["page_number"]
        page_text = page_doc["text"]

        if enable_cleaning and HAS_CLEANERS:
            page_text = preprocess_for_embedding(page_text, is_markdown=False, is_pdf=True)

        if not page_text or not page_text.strip():
            continue

        page_chunk_list = create_chunks(
            page_text,
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            file_type="pdf",
        )
        if not page_chunk_list:
            logger.warning(f"No chunks created for PDF page {page_number}")
            continue

        for local_idx, chunk in enumerate(page_chunk_list):
            out.append({
                "page_content": chunk,
                "metadata": {
                    "file_type": "pdf",
                    "chunk_index": global_idx,
                    "page_chunk_index": local_idx,
                    "chunk_strategy": "pdf",
                    "page_number": page_number,
                    "source": filename or str(file_path),
                    "file_path": std_source_path,
                    "file_name": Path(filename or str(file_path)).name,
                    # Standard path never OCRs — consistent native-text provenance.
                    "extraction_method": "native_text",
                    "ocr_used": False,
                    "ocr_provider": None,
                    "from_image": False,
                    "table_detected": False,
                    "table_index": None,
                    "confidence": None,
                    "language": None,
                },
            })
            global_idx += 1

    return out


def load_document(
    file_path: Union[str, Path, io.BytesIO],
    filename: str = None,
    enable_cleaning: bool = True,
    enable_chunking: bool = True,
    chunk_size: int = None,
    chunk_overlap: int = None
) -> List[Dict]:

    try:
        is_stream = isinstance(file_path, io.BytesIO)
        if not is_stream:
            file_path = resolve_shared_server_path(str(file_path))
            file_path_obj = Path(file_path)

            logger.info(f"Verifying file: {file_path_obj}")

            if not file_path_obj.exists():
                logger.error(f"File does not exist: {file_path_obj}")
                return []

            logger.info(f"File verified ({file_path_obj.stat().st_size} bytes)")
        else:
            file_path_obj = file_path

        ext = get_file_extension(file_path_obj, filename)

        if not ext:
            mime_type, _ = mimetypes.guess_type(str(file_path))
            if mime_type:
                ext = mime_type.split("/")[-1]

        ext = (ext or "").lower()

        logger.info(f"Detected file type: {ext or 'unknown'}")

        if enable_chunking and HAS_CHUNKER and settings.ENABLE_LAYOUT_AWARE_CHUNKING and ext in ("pdf", "pptx", "ppt"):
            structured_chunks = _try_layout_aware_chunking(file_path_obj, ext, filename or str(file_path))
            if structured_chunks is not None:
                return structured_chunks
            logger.info("Layout-aware path skipped/failed, falling back to standard chunking")

        # Plain function refs (not lambdas) so the parser name can be recorded
        # in each chunk's extraction metadata via loader.__name__.
        loaders = {
            "pdf": load_pdf,
            "docx": load_word_document,
            "doc": load_word_document,
            "xlsx": load_excel_file,
            "xls": load_excel_file,
            "pptx": load_powerpoint,
            "ppt": load_powerpoint,
            "txt": load_text_file,
            "text": load_text_file,
            "md": load_markdown,
            "markdown": load_markdown,
            "csv": load_csv,
            "json": load_json,
        }

        # ── Page-aware standard path for PDFs ──────────────────────────────
        # The flat path joins every page into one string, so page boundaries are
        # lost and chunks can't carry a page_number. For PDFs we instead extract
        # one text block per page and chunk each page separately, stamping
        # page_number onto every chunk. Generic — no per-document assumptions.
        if (
            enable_chunking and HAS_CHUNKER and ext == "pdf"
            and not isinstance(file_path_obj, io.BytesIO)
        ):
            page_docs = load_pdf_pages(file_path_obj)
            if page_docs:
                page_chunks = _chunk_pdf_pages_with_numbers(
                    page_docs,
                    filename=filename,
                    file_path=file_path,
                    file_path_obj=file_path_obj,
                    enable_cleaning=enable_cleaning,
                    chunk_size=chunk_size,
                    chunk_overlap=chunk_overlap,
                )
                if page_chunks:
                    logger.info(
                        f"Created {len(page_chunks)} chunks "
                        f"(page-aware, {len(page_docs)} pages)"
                    )
                    return page_chunks
                logger.info("Page-aware PDF chunking produced nothing, falling back to flat")

        loader = loaders.get(ext, load_text_file)

        content = loader(file_path_obj)

        if not content:
            logger.error("No content extracted")
            return []

        logger.info(f"Extracted {len(content)} characters")

        # Markdown structure stats must be computed on the RAW content, before
        # preprocess_for_embedding strips/normalizes anything.
        md_stats = (
            markdown_content_stats(content)
            if ext in ("md", "markdown") and isinstance(content, str)
            else None
        )

        is_markdown = ext in ["md", "markdown"]
        is_pdf = ext == "pdf"

        # xlsx/xls/csv loaders return a structured dict (headers + rows) that
        # goes straight to the row-aware chunker — text cleaning doesn't apply.
        if enable_cleaning and HAS_CLEANERS and ext not in ['xlsx', 'xls', 'csv']:
            content = preprocess_for_embedding(
                content,
                is_markdown=is_markdown,
                is_pdf=is_pdf,
                # Word tables arrive as <doc-table> HTML; keep them intact so
                # chunk_structured_content can treat them as atomic units.
                preserve_doc_tables=ext in ('docx', 'doc'),
            )
        
        if not enable_chunking or not HAS_CHUNKER:
            return [{
                "page_content": content,
                "metadata": {"file_type": ext}
            }]

        file_type_map = {
            "pdf": "pdf",
            "docx": "word",
            "doc": "word",
            "xlsx": "excel",
            "xls": "excel",
            # CSV is tabular data — same header-aware row chunker as Excel.
            "csv": "excel",
            "md": "markdown",
            "markdown": "markdown",
            "txt": "text",
            "text": "text",
        }

        file_type = file_type_map.get(ext)

        logger.info(f"Chunking using strategy: {file_type or 'default'}")

        chunk_list = create_chunks(
            content,
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            file_type=file_type
        )

        if not chunk_list:
            logger.warning("No chunks created")
            return []

        std_source_path = str(file_path_obj) if not isinstance(file_path_obj, io.BytesIO) else None
        chunks = [
            {
                "page_content": chunk,
                "metadata": {
                    "file_type": ext,
                    "chunk_index": idx,
                    "chunk_strategy": file_type or "default",
                    "source": filename or str(file_path),
                    "file_path": std_source_path,
                    "file_name": Path(filename or str(file_path)).name,
                    # Standard path never OCRs — stamp native-text provenance so
                    # every chunk carries consistent extraction metadata.
                    "extraction_method": "native_text",
                    "ocr_used": False,
                    "ocr_provider": None,
                    "from_image": False,
                    "table_detected": False,
                    "table_index": None,
                    "confidence": None,
                    "language": None,
                    "parser": getattr(loader, "__name__", "unknown"),
                    "mime_type": mimetypes.guess_type(filename or str(file_path))[0],
                    "extraction_status": "success",
                    "warnings": [],
                    # Markdown-only document structure stats (heading_count,
                    # link_count, code_block_count, table_count).
                    **(md_stats or {}),
                }
            }
            for idx, chunk in enumerate(chunk_list)
        ]

        logger.info(f"Created {len(chunks)} chunks")
        return chunks

    except Exception as e:
        logger.error(f"load_document failed: {e}")
        import traceback
        logger.error(traceback.format_exc())
        return []


def load_multiple_documents(
    file_paths: List[Union[str, Path]],
    enable_cleaning: bool = True,
    enable_chunking: bool = True
) -> List[Dict[str, str]]:
    """
    PRODUCTION: Load multiple documents with cleaning and chunking

    Args:
        file_paths: List of file paths
        enable_cleaning: Apply text cleaning
        enable_chunking: Apply file-type-specific chunking

    Returns:
        List of dicts with filename, filepath, chunks, and metadata
    """
    documents = []
    for file_path in file_paths:
        try:
            chunks = load_document(
                file_path,
                enable_cleaning=enable_cleaning,
                enable_chunking=enable_chunking
            )
            if chunks:
                documents.append({
                    'filename': Path(file_path).name,
                    'filepath': str(file_path),
                    'chunks': chunks,
                    'num_chunks': len(chunks),
                    'total_length': sum(len(chunk['page_content']) for chunk in chunks)
                })
                logger.info(f"Loaded: {Path(file_path).name} ({len(chunks)} chunks)")
            else:
                logger.warning(f"Empty content: {Path(file_path).name}")
        except Exception as e:
            logger.error(f"Failed to load {Path(file_path).name}: {e}")

    total_chunks = sum(doc['num_chunks'] for doc in documents)
    logger.info(f"\nLoaded {len(documents)}/{len(file_paths)} documents ({total_chunks} total chunks)")
    return documents


def get_supported_formats() -> Dict[str, List[str]]:
    """Get list of supported file formats"""
    formats = {
        # PyMuPDF (fitz) is the primary PDF parser and a hard dependency;
        # PyPDF2/pdfplumber are optional extras.
        'pdf': ['pdf'],
        'word': ['docx', 'doc'] if HAS_DOCX else [],
        'excel': ['xlsx', 'xls'] if HAS_OPENPYXL else [],
        'powerpoint': ['pptx', 'ppt'] if HAS_PPTX else [],
        'text': ['txt', 'text'],
        'markdown': ['md', 'markdown'],
        'data': ['csv', 'json'],
        'web': ['html', 'htm']
    }

    return {k: v for k, v in formats.items() if v}


def print_supported_formats():
    """Print all supported file formats and features"""
    formats = get_supported_formats()
    print("\n" + "="*60)
    print("SUPPORTED FILE FORMATS")
    print("="*60)
    for category, extensions in formats.items():
        print(f"\n{category.upper()}:")
        print(f"  {', '.join(extensions)}")
    print("\n" + "="*60)

    missing = []
    if not (HAS_PYPDF2 or HAS_PDFPLUMBER):
        missing.append("PyPDF2 or pdfplumber (optional PDF fallbacks; PyMuPDF is primary)")
    if not HAS_DOCX:
        missing.append("python-docx (for Word)")
    if not HAS_OPENPYXL:
        missing.append("openpyxl (for Excel)")
    if not HAS_PPTX:
        missing.append("python-pptx (for PowerPoint)")

    if missing:
        print("\nInstall optional dependencies for more formats:")
        for lib in missing:
            print(f"  - {lib}")

    print("\n" + "="*60)
    print("INTEGRATED FEATURES")
    print("="*60)
    print(f"\nText Cleaning: {'Enabled' if HAS_CLEANERS else 'Disabled (cleaners.py missing)'}")
    print(f"Smart Chunking: {'Enabled' if HAS_CHUNKER else 'Disabled (chunker.py missing)'}")

    if HAS_CHUNKER:
        print("\nCHUNKING STRATEGIES:")
        print("  PDF:      768 tokens, 150 overlap, table-preserving")
        print("  Word:     512 tokens, 100 overlap, recursive")
        print("  Excel:    15 rows/chunk, 2 row overlap, semantic")
        print("  Text:     512 tokens, 100 overlap, recursive")
        print("  Markdown: 1024 tokens, 100 overlap, header-aware")

    print("\n" + "="*60)